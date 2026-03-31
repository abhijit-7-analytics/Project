"""
SalesDB — 15-Slide Professional Presentation
Theme: White / Clean | Font: Times New Roman
Title: <40pt | Body: <26pt
Run:   python ppt.py
Output: SalesDB_Presentation.pptx + ppt_images/
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, Rectangle
import numpy as np, os

IMG = "ppt_images"
os.makedirs(IMG, exist_ok=True)
FONT = 'Times New Roman'

# ═══ PPT COLORS (White Theme) ═══
BG      = RGBColor(255,255,255)
LIGHT   = RGBColor(248,249,250)
TITLE_C = RGBColor(26,26,46)
BODY_C  = RGBColor(51,65,85)
ACCENT  = RGBColor(124,58,237)
BLUE    = RGBColor(37,99,235)
GREEN   = RGBColor(16,185,129)
PINK    = RGBColor(236,72,153)
RED     = RGBColor(239,68,68)
ORANGE  = RGBColor(249,115,22)
CYAN    = RGBColor(6,182,212)
YELLOW  = RGBColor(245,158,11)
BORDER  = RGBColor(229,231,235)
MUTED   = RGBColor(107,114,128)
WHITE   = RGBColor(255,255,255)
BLACK   = RGBColor(0,0,0)

# ═══ MATPLOTLIB COLORS ═══
W='#ffffff'; WC='#f8f9fa'; WB='#e5e7eb'
WP='#7c3aed'; WBL='#2563eb'; WG='#10b981'
WPI='#ec4899'; WR='#ef4444'; WO='#f97316'
WCY='#06b6d4'; WY='#f59e0b'; WM='#6b7280'
WT='#1e293b'; WT2='#334155'
# Dark (for app screenshots)
DB='#0f172a'; DC='#1e293b'; DBR='#334155'
DP='#ec4899'; DBL='#3b82f6'; DPU='#a855f7'
DG='#3ddc84'; DM='#94a3b8'; DR='#f43f5e'
DY='#fbbf24'; DCY='#06b6d4'


# ══════════════════════════════════════════
#  IMAGE GENERATORS — WHITE BG DIAGRAMS
# ══════════════════════════════════════════

def gen_tech_stack():
    fig, ax = plt.subplots(figsize=(10, 5))
    ax.set_xlim(0, 10); ax.set_ylim(0, 6); ax.axis('off')
    layers = [
        (0.5, 4.5, 9, 0.85, 'Frontend :  HTML5  |  CSS3  |  JavaScript ES6  |  Chart.js', WBL),
        (0.5, 3.3, 9, 0.85, 'Backend  :  Python 3.12  |  Flask 3.0  |  Flask-CORS', WP),
        (0.5, 2.1, 9, 0.85, 'Security :  Flask Sessions  |  @login_required  |  Cookies', WO),
        (0.5, 0.9, 9, 0.85, 'Database :  PostgreSQL 16  |  Star Schema  |  5 Tables', WG),
    ]
    for x, y, w, h, label, color in layers:
        ax.add_patch(FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.08",
                     facecolor=color, edgecolor='white', lw=2.5, alpha=0.9))
        tc = 'white' if color != WY else 'black'
        ax.text(x + w / 2, y + h / 2, label, ha='center', va='center',
                fontsize=11, fontweight='bold', color=tc, fontfamily='serif')
    for i in range(3):
        y = 4.5 - i * 1.2
        ax.annotate('', xy=(5, y), xytext=(5, y + 0.35),
                    arrowprops=dict(arrowstyle='<->', color=WM, lw=2))
    p = os.path.join(IMG, "tech_stack.png")
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor=W); plt.close(); return p


def gen_er_diagram():
    fig, ax = plt.subplots(figsize=(14, 9))
    ax.set_xlim(0, 14); ax.set_ylim(0, 9); ax.axis('off')

    tables = {
        'customer': {
            'x': 0.3, 'y': 5, 'w': 3.5, 'h': 3.5, 'color': WBL,
            'title': 'CUSTOMER_DIM',
            'fields': [
                ('PK', 'customer_id', 'SERIAL'),
                ('', 'first_name', 'VARCHAR(50)'),
                ('', 'last_name', 'VARCHAR(50)'),
                ('', 'city', 'VARCHAR(50)'),
                ('', 'mobile_no', 'VARCHAR(20)'),
                ('', 'email', 'VARCHAR(100)'),
                ('', 'region', 'VARCHAR(20)'),
                ('', 'member_type', 'VARCHAR(20)'),
            ]
        },
        'product': {
            'x': 10.2, 'y': 5, 'w': 3.5, 'h': 2.8, 'color': WP,
            'title': 'PRODUCT_DIM',
            'fields': [
                ('PK', 'product_id', 'SERIAL'),
                ('', 'product_name', 'VARCHAR(100)'),
                ('', 'category', 'VARCHAR(50)'),
                ('', 'unit_price', 'NUMERIC(10,2)'),
                ('', 'stock_qty', 'INTEGER'),
            ]
        },
        'sales': {
            'x': 5.2, 'y': 5.5, 'w': 3.5, 'h': 3, 'color': WPI,
            'title': 'SALES_FACT',
            'fields': [
                ('PK', 'sale_id', 'SERIAL'),
                ('FK', 'customer_id', 'INTEGER'),
                ('FK', 'product_id', 'INTEGER'),
                ('', 'sale_date', 'DATE'),
                ('', 'quantity', 'INTEGER'),
                ('', 'sale_amount', 'NUMERIC(10,2)'),
            ]
        },
        'invoice': {
            'x': 1.5, 'y': 0.3, 'w': 3.8, 'h': 3.8, 'color': WCY,
            'title': 'INVOICE_FACT',
            'fields': [
                ('PK', 'invoice_id', 'SERIAL'),
                ('FK', 'sale_id', 'INTEGER'),
                ('', 'invoice_no', 'VARCHAR(30)'),
                ('', 'subtotal', 'NUMERIC(10,2)'),
                ('', 'tax_rate', 'NUMERIC(5,2)'),
                ('', 'tax_amount', 'NUMERIC(10,2)'),
                ('', 'grand_total', 'NUMERIC(10,2)'),
                ('', 'status', 'VARCHAR(20)'),
                ('', 'created_at', 'TIMESTAMP'),
            ]
        },
        'stock': {
            'x': 9, 'y': 0.3, 'w': 3.5, 'h': 3, 'color': WO,
            'title': 'STOCK_HISTORY',
            'fields': [
                ('PK', 'history_id', 'SERIAL'),
                ('FK', 'product_id', 'INTEGER'),
                ('', 'change_type', 'VARCHAR(20)'),
                ('', 'qty_change', 'INTEGER'),
                ('', 'new_stock', 'INTEGER'),
                ('', 'created_at', 'TIMESTAMP'),
            ]
        },
    }

    for key, t in tables.items():
        hdr = 0.5
        # Table body
        ax.add_patch(FancyBboxPatch((t['x'], t['y']), t['w'], t['h'],
                     boxstyle="round,pad=0.06", facecolor='#fafbfc',
                     edgecolor=t['color'], lw=2.5))
        # Header
        ax.add_patch(FancyBboxPatch((t['x'], t['y'] + t['h'] - hdr), t['w'], hdr,
                     boxstyle="round,pad=0.06", facecolor=t['color'],
                     edgecolor=t['color'], lw=2.5))
        ax.text(t['x'] + t['w'] / 2, t['y'] + t['h'] - hdr / 2, t['title'],
                ha='center', va='center', fontsize=9, fontweight='bold',
                color='white', fontfamily='monospace')
        # Fields
        for i, (pk, fname, ftype) in enumerate(t['fields']):
            fy = t['y'] + t['h'] - hdr - 0.35 - i * 0.32
            fc = WO if pk == 'PK' else WR if pk == 'FK' else WT2
            marker = '🔑 ' if pk == 'PK' else '🔗 ' if pk == 'FK' else '    '
            ax.text(t['x'] + 0.15, fy, f"{marker}{fname}", fontsize=7.5,
                    color=fc, fontfamily='monospace', va='center')
            ax.text(t['x'] + t['w'] - 0.15, fy, ftype, fontsize=6.5,
                    color=WM, fontfamily='monospace', va='center', ha='right')

    # Relationships
    # customer → sales (1:N)
    ax.annotate('', xy=(5.2, 7), xytext=(3.8, 7),
                arrowprops=dict(arrowstyle='->', color=WBL, lw=2.5))
    ax.text(4.5, 7.25, '1 : N', fontsize=10, fontweight='bold', color=WBL, ha='center')

    # product → sales (1:N)
    ax.annotate('', xy=(8.7, 7), xytext=(10.2, 7),
                arrowprops=dict(arrowstyle='->', color=WP, lw=2.5))
    ax.text(9.5, 7.25, '1 : N', fontsize=10, fontweight='bold', color=WP, ha='center')

    # sales → invoice (1:1)
    ax.annotate('', xy=(4, 4.1), xytext=(5.8, 5.5),
                arrowprops=dict(arrowstyle='->', color=WCY, lw=2.5))
    ax.text(4.5, 5, '1 : 1', fontsize=10, fontweight='bold', color=WCY, ha='center')

    # product → stock_history (1:N)
    ax.annotate('', xy=(10.5, 3.3), xytext=(11.5, 5),
                arrowprops=dict(arrowstyle='->', color=WO, lw=2.5))
    ax.text(11.5, 4.2, '1 : N', fontsize=10, fontweight='bold', color=WO, ha='center')

    p = os.path.join(IMG, "er_diagram.png")
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor=W); plt.close(); return p


def gen_dfd_level0():
    fig, ax = plt.subplots(figsize=(10, 5))
    ax.set_xlim(0, 10); ax.set_ylim(0, 5); ax.axis('off')

    # External entity
    ax.add_patch(FancyBboxPatch((0.3, 1.5), 2, 2, boxstyle="round,pad=0.1",
                 facecolor=WBL, edgecolor='white', lw=2))
    ax.text(1.3, 2.5, 'Admin\nUser', ha='center', va='center',
            fontsize=13, fontweight='bold', color='white')

    # Process
    c = plt.Circle((5, 2.5), 1.3, facecolor=WP, edgecolor='white', lw=2.5)
    ax.add_patch(c)
    ax.text(5, 2.5, 'SalesDB\nSystem', ha='center', va='center',
            fontsize=13, fontweight='bold', color='white')

    # Data store
    ax.add_patch(FancyBboxPatch((7.5, 1.5), 2.2, 2, boxstyle="round,pad=0.1",
                 facecolor=WG, edgecolor='white', lw=2))
    ax.text(8.6, 2.5, 'PostgreSQL\nDatabase', ha='center', va='center',
            fontsize=11, fontweight='bold', color='white')

    # Arrows
    ax.annotate('', xy=(3.7, 2.8), xytext=(2.3, 2.8),
                arrowprops=dict(arrowstyle='->', color=WT, lw=2))
    ax.text(2.8, 3.5, 'Login, CRUD\nRequests', fontsize=9, ha='center', color=WT2)
    ax.annotate('', xy=(2.3, 2.2), xytext=(3.7, 2.2),
                arrowprops=dict(arrowstyle='->', color=WM, lw=2))
    ax.text(2.8, 1.1, 'Dashboard,\nReports, PDFs', fontsize=9, ha='center', color=WM)
    ax.annotate('', xy=(7.5, 2.8), xytext=(6.3, 2.8),
                arrowprops=dict(arrowstyle='->', color=WT, lw=2))
    ax.text(6.8, 3.5, 'SQL\nQueries', fontsize=9, ha='center', color=WT2)
    ax.annotate('', xy=(6.3, 2.2), xytext=(7.5, 2.2),
                arrowprops=dict(arrowstyle='->', color=WM, lw=2))
    ax.text(6.8, 1.1, 'Result\nSets', fontsize=9, ha='center', color=WM)

    p = os.path.join(IMG, "dfd_level0.png")
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor=W); plt.close(); return p


def gen_dfd_level1():
    fig, ax = plt.subplots(figsize=(13, 9))
    ax.set_xlim(0, 13); ax.set_ylim(0, 10); ax.axis('off')

    # User entity
    ax.add_patch(FancyBboxPatch((0.2, 4), 1.8, 1.5, boxstyle="round,pad=0.1",
                 facecolor=WBL, edgecolor='white', lw=2))
    ax.text(1.1, 4.75, 'Admin\nUser', ha='center', va='center',
            fontsize=10, fontweight='bold', color='white')

    # Processes
    procs = [
        (4, 8.5, '0.0', 'Authenticate', WY),
        (4, 7, '1.0', 'Manage\nCustomers', WPI),
        (4, 5.5, '2.0', 'Manage\nProducts', WP),
        (4, 4, '3.0', 'Manage\nSales', WR),
        (4, 2.5, '4.0', 'Generate\nInvoices', WCY),
        (4, 1, '5.0', 'Manage\nStocks', WO),
        (8, 1, '6.0', 'Generate\nReports', WG),
    ]
    for x, y, num, label, color in procs:
        c = plt.Circle((x + 1, y + 0.5), 0.62, facecolor=color,
                       edgecolor='white', lw=2)
        ax.add_patch(c)
        tc = 'white' if color != WY else 'black'
        ax.text(x + 1, y + 0.7, num, ha='center', va='center',
                fontsize=7, fontweight='bold', color=tc)
        ax.text(x + 1, y + 0.3, label, ha='center', va='center',
                fontsize=6, color=tc)

    # Data stores
    stores = [
        (8.5, 8.3, 'D1', 'SESSION', WY),
        (8.5, 7.1, 'D2', 'CUSTOMER_DIM', WBL),
        (8.5, 5.9, 'D3', 'PRODUCT_DIM', WP),
        (8.5, 4.7, 'D4', 'SALES_FACT', WPI),
        (8.5, 3.5, 'D5', 'INVOICE_FACT', WCY),
        (8.5, 2.3, 'D6', 'STOCK_HISTORY', WO),
    ]
    for x, y, did, label, color in stores:
        ax.add_patch(Rectangle((x, y), 4, 0.65, facecolor=WC,
                              edgecolor=color, lw=2))
        ax.plot([x + 0.6, x + 0.6], [y, y + 0.65], color=color, lw=2)
        ax.text(x + 0.3, y + 0.32, did, ha='center', va='center',
                fontsize=7, fontweight='bold', color=color)
        ax.text(x + 0.6 + 1.7, y + 0.32, label, ha='center', va='center',
                fontsize=7.5, color=WT, fontfamily='monospace')

    # Arrows from user to processes
    for proc in procs[:6]:
        ax.annotate('', xy=(proc[0] + 0.38, proc[1] + 0.5),
                    xytext=(2, 4.75),
                    arrowprops=dict(arrowstyle='->', color=WM, lw=1))
    # Arrows from processes to stores
    conns = [(5.62, 9, 8.5, 8.62), (5.62, 7.5, 8.5, 7.42),
             (5.62, 6, 8.5, 6.22), (5.62, 4.5, 8.5, 5.02),
             (5.62, 3, 8.5, 3.82), (5.62, 1.5, 8.5, 2.62),
             (9.62, 1.5, 10.5, 2.3)]
    for x1, y1, x2, y2 in conns:
        ax.annotate('', xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle='->', color=WM, lw=1))

    p = os.path.join(IMG, "dfd_level1.png")
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor=W); plt.close(); return p


def gen_auth_flow():
    fig, ax = plt.subplots(figsize=(10, 5))
    ax.set_xlim(0, 10); ax.set_ylim(0, 6); ax.axis('off')

    steps = [
        (0.3, 4, 2, 1.2, 'User clicks\nprotected page', WBL),
        (2.8, 4, 2, 1.2, 'Login overlay\nappears', WPI),
        (5.3, 4, 2, 1.2, 'POST /api/login\nvalidate creds', WP),
        (5.3, 1.5, 2, 1.2, 'Session created\n(30 min expiry)', WG),
        (2.8, 1.5, 2, 1.2, 'UI unlocks\n✓ Admin badge', WY),
        (0.3, 1.5, 2, 1.2, 'Access CRUD\npages', WBL),
    ]
    for x, y, w, h, label, color in steps:
        ax.add_patch(FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.1",
                     facecolor=color, edgecolor='white', lw=2))
        tc = 'white' if color != WY else 'black'
        ax.text(x + w / 2, y + h / 2, label, ha='center', va='center',
                fontsize=9, fontweight='bold', color=tc)

    arrows = [(2.3, 4.6, 2.8, 4.6), (4.8, 4.6, 5.3, 4.6),
              (6.3, 4, 6.3, 2.7), (5.3, 2.1, 4.8, 2.1), (2.8, 2.1, 2.3, 2.1)]
    for x1, y1, x2, y2 in arrows:
        ax.annotate('', xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle='->', color=WT2, lw=2.5))

    # Session expiry
    ax.add_patch(FancyBboxPatch((8, 2.5), 1.6, 0.9, boxstyle="round,pad=0.08",
                 facecolor=WR, edgecolor='white', lw=2))
    ax.text(8.8, 2.95, 'Session\nExpires', ha='center', va='center',
            fontsize=8, fontweight='bold', color='white')
    ax.annotate('', xy=(8, 2.95), xytext=(7.3, 2.1),
                arrowprops=dict(arrowstyle='->', color=WR, lw=1.5, ls='dashed'))

    p = os.path.join(IMG, "auth_flow.png")
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor=W); plt.close(); return p


def gen_result_charts():
    fig, axes = plt.subplots(1, 3, figsize=(14, 4.5))
    cats = ['Electronics', 'Home Goods', 'Apparel', 'Other']
    rev = [45200, 28100, 17800, 8900]
    colors = [WPI, WBL, WP, WO]

    # Bar chart
    axes[0].bar(cats, rev, color=colors, edgecolor='white', lw=1, width=0.55)
    axes[0].set_title('Revenue by Category', fontsize=12, fontweight='bold',
                      color=WT, pad=10, fontfamily='serif')
    axes[0].set_facecolor(WC)
    axes[0].tick_params(colors=WM, labelsize=8)
    for sp in ['top', 'right']:
        axes[0].spines[sp].set_visible(False)
    for sp in ['bottom', 'left']:
        axes[0].spines[sp].set_color(WB)
    axes[0].yaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f'${x / 1000:.0f}K'))

    # Doughnut
    axes[1].pie(rev, labels=cats, autopct='%1.1f%%', colors=colors,
                textprops={'fontsize': 8, 'color': WT}, startangle=90,
                pctdistance=0.78, wedgeprops=dict(width=0.4, edgecolor=W, linewidth=2))
    axes[1].set_title('Sales Distribution', fontsize=12, fontweight='bold',
                      color=WT, pad=10, fontfamily='serif')

    # Pie
    axes[2].pie([42, 58], labels=['Gold  42%', 'Regular  58%'],
                colors=[WY, '#94a3b8'], textprops={'fontsize': 10, 'color': WT},
                startangle=90, wedgeprops=dict(edgecolor=W, linewidth=2))
    axes[2].set_title('Customer Segments', fontsize=12, fontweight='bold',
                      color=WT, pad=10, fontfamily='serif')

    fig.patch.set_facecolor(W)
    plt.tight_layout(pad=2)
    p = os.path.join(IMG, "result_charts.png")
    plt.savefig(p, dpi=200, bbox_inches='tight', facecolor=W); plt.close(); return p


# ══════════════════════════════════════════
#  SCREENSHOT GENERATORS (Dark BG — actual app)
# ══════════════════════════════════════════

def _sidebar(ax, active_idx, xlim=12, ylim=8):
    ax.set_xlim(0, xlim); ax.set_ylim(0, ylim); ax.axis('off')
    ax.add_patch(Rectangle((0, 0), xlim, ylim, facecolor=DB, edgecolor=DBR, lw=2))
    ax.add_patch(Rectangle((0, 0), 2.2, ylim, facecolor=DC, edgecolor=DBR, lw=1))
    ax.text(0.3, ylim - 0.5, 'SalesDB', fontsize=10, fontweight='bold', color=DP)
    ax.text(0.3, ylim - 0.85, 'Analytics', fontsize=6, color=DM)
    nav = ['▦ Dashboard', '◈ Sales', '◉ Customers', '◇ Products',
           '▤ Invoices', '▣ Stocks', '▥ Reports']
    for i, item in enumerate(nav):
        y = ylim - 1.5 - i * 0.45
        fc = DP if i == active_idx else DM
        ax.text(0.3, y, item, fontsize=6, color=fc)
        if i != active_idx and i > 0:
            ax.text(2, y, '🔒', fontsize=4, color=DM)
    ax.add_patch(Rectangle((2.2, ylim - 0.8), xlim - 2.2, 0.8,
                           facecolor=DC, edgecolor=DBR, lw=1))


def gen_ss_dashboard():
    fig, ax = plt.subplots(figsize=(12, 7))
    _sidebar(ax, 0)
    ax.text(2.6, 7.25, 'Dashboard', fontsize=11, fontweight='bold', color='white')
    ax.add_patch(FancyBboxPatch((9.5, 7.1), 1.5, 0.45, boxstyle="round,pad=0.04",
                 facecolor=DY, edgecolor=DY, lw=1))
    ax.text(10.25, 7.32, '🔒 Login', ha='center', fontsize=7, fontweight='bold', color='black')
    kpis = [('$100K', DP), ('245', DBL), ('$408', DPU), ('48', DR), ('12', DG)]
    for i, (val, color) in enumerate(kpis):
        x = 2.6 + i * 1.85
        ax.add_patch(FancyBboxPatch((x, 5.2), 1.65, 1.3, boxstyle="round,pad=0.06",
                     facecolor=DC, edgecolor=DBR, lw=1))
        ax.add_patch(Rectangle((x, 6.47), 1.65, 0.03, facecolor=color))
        ax.text(x + 0.15, 5.85, val, fontsize=14, fontweight='bold', color='white')
    for i in range(3):
        x = 2.6 + i * 3.1
        ax.add_patch(FancyBboxPatch((x, 0.5), 2.9, 4.3, boxstyle="round,pad=0.06",
                     facecolor=DC, edgecolor=DBR, lw=1))
        if i == 0:
            for j, h in enumerate([2.5, 1.8, 1.2, 0.6]):
                ax.add_patch(Rectangle((x + 0.3 + j * 0.6, 0.8), 0.4, h,
                            facecolor=[DP, DBL, DPU, DR][j]))
        elif i == 1:
            c1 = plt.Circle((x + 1.45, 2.3), 1, facecolor=DP, edgecolor=DB, lw=0)
            ax.add_patch(c1)
            c2 = plt.matplotlib.patches.Wedge((x + 1.45, 2.3), 1, 90, 220,
                 facecolor=DBL, edgecolor=DB, lw=0)
            ax.add_patch(c2)
            c3 = plt.Circle((x + 1.45, 2.3), 0.5, facecolor=DC); ax.add_patch(c3)
        else:
            c1 = plt.Circle((x + 1.45, 2.3), 1, facecolor=DY, edgecolor=DB, lw=0)
            ax.add_patch(c1)
            c2 = plt.matplotlib.patches.Wedge((x + 1.45, 2.3), 1, 0, 210,
                 facecolor='#475569', edgecolor=DB, lw=0)
            ax.add_patch(c2)
    p = os.path.join(IMG, "ss_dashboard.png")
    plt.savefig(p, dpi=180, bbox_inches='tight', facecolor=DB); plt.close(); return p


def gen_ss_login():
    fig, ax = plt.subplots(figsize=(10, 7))
    ax.set_xlim(0, 10); ax.set_ylim(0, 7); ax.axis('off')
    ax.add_patch(Rectangle((0, 0), 10, 7, facecolor='#0a0e1a', alpha=0.95))
    cx, cy, cw, ch = 2.5, 1.2, 5, 4.6
    ax.add_patch(FancyBboxPatch((cx, cy), cw, ch, boxstyle="round,pad=0.15",
                 facecolor=DC, edgecolor=DBR, lw=2))
    ax.add_patch(Rectangle((cx, cy + ch - 0.08), cw, 0.08, facecolor=DY))
    ax.text(5, 5.2, 'SalesDB', ha='center', fontsize=24, fontweight='bold', color='white')
    ax.text(5, 4.8, 'ADMIN AUTHENTICATION', ha='center', fontsize=8,
            color=DM, fontweight='bold')
    ax.add_patch(FancyBboxPatch((3, 3.8), 4, 0.5, boxstyle="round,pad=0.05",
                 facecolor=DB, edgecolor=DBR, lw=1.5))
    ax.text(3.15, 4.05, 'admin', fontsize=11, color='white')
    ax.add_patch(FancyBboxPatch((3, 3.1), 4, 0.5, boxstyle="round,pad=0.05",
                 facecolor=DB, edgecolor=DBR, lw=1.5))
    ax.text(3.15, 3.35, '• • • • • • • •', fontsize=11, color=DM)
    ax.add_patch(FancyBboxPatch((3, 2.2), 4, 0.55, boxstyle="round,pad=0.05",
                 facecolor=DY, edgecolor=DY, lw=1.5))
    ax.text(5, 2.47, 'LOGIN', ha='center', fontsize=12, fontweight='bold', color='black')
    ax.text(5, 1.7, 'Session expires after 30 minutes', ha='center',
            fontsize=7, color=DM)
    p = os.path.join(IMG, "ss_login.png")
    plt.savefig(p, dpi=180, bbox_inches='tight', facecolor='#0a0e1a'); plt.close(); return p


def gen_ss_sales():
    fig, ax = plt.subplots(figsize=(12, 7))
    _sidebar(ax, 1)
    ax.text(2.6, 7.25, 'Sales', fontsize=11, fontweight='bold', color='white')
    ax.add_patch(FancyBboxPatch((2.6, 4.5), 9, 2.3, boxstyle="round,pad=0.06",
                 facecolor=DC, edgecolor=DBR, lw=1))
    ax.add_patch(Rectangle((2.6, 6.5), 9, 0.3, facecolor=DBR))
    ax.text(2.8, 6.55, 'Add New Sale', fontsize=8, fontweight='bold', color='white')
    ax.add_patch(FancyBboxPatch((2.6, 0.3), 9, 3.8, boxstyle="round,pad=0.06",
                 facecolor=DC, edgecolor=DBR, lw=1))
    ax.add_patch(Rectangle((2.6, 3.8), 9, 0.3, facecolor=DBR))
    ax.text(2.8, 3.85, 'All Sales', fontsize=8, fontweight='bold', color='white')
    headers = ['#', 'Customer', 'Product', 'Cat', 'Date', 'Qty', 'Price', 'Total']
    for i, h in enumerate(headers):
        ax.text(2.7 + i * 1.1, 3.45, h, fontsize=5.5, color=DM, fontweight='bold')
    p = os.path.join(IMG, "ss_sales.png")
    plt.savefig(p, dpi=180, bbox_inches='tight', facecolor=DB); plt.close(); return p


def gen_ss_products():
    fig, ax = plt.subplots(figsize=(12, 7))
    _sidebar(ax, 3)
    ax.text(2.6, 7.25, 'Products', fontsize=11, fontweight='bold', color='white')
    ax.add_patch(FancyBboxPatch((2.6, 5), 9, 1.8, boxstyle="round,pad=0.06",
                 facecolor=DC, edgecolor=DBR, lw=1))
    ax.add_patch(Rectangle((2.6, 6.5), 9, 0.3, facecolor=DBR))
    ax.text(2.8, 6.55, 'Add Product', fontsize=8, fontweight='bold', color='white')
    ax.add_patch(FancyBboxPatch((2.6, 0.3), 9, 4.3, boxstyle="round,pad=0.06",
                 facecolor=DC, edgecolor=DBR, lw=1))
    ax.add_patch(Rectangle((2.6, 4.3), 9, 0.3, facecolor=DBR))
    ax.text(2.8, 4.35, 'All Products', fontsize=8, fontweight='bold', color='white')
    p = os.path.join(IMG, "ss_products.png")
    plt.savefig(p, dpi=180, bbox_inches='tight', facecolor=DB); plt.close(); return p


def gen_ss_invoices():
    fig, ax = plt.subplots(figsize=(12, 7))
    _sidebar(ax, 4)
    ax.text(2.6, 7.25, 'Invoices', fontsize=11, fontweight='bold', color='white')
    ax.add_patch(FancyBboxPatch((2.6, 3.5), 9, 3.2, boxstyle="round,pad=0.06",
                 facecolor=DC, edgecolor=DBR, lw=1))
    ax.add_patch(Rectangle((2.6, 6.4), 9, 0.3, facecolor=DPU))
    ax.text(2.8, 6.47, '📄 Invoice Builder', fontsize=8, fontweight='bold', color='white')
    ax.text(10.5, 6.47, 'INV-202401-12345', fontsize=7, fontweight='bold',
            color=DP, ha='right')
    ax.add_patch(FancyBboxPatch((2.6, 0.3), 9, 2.8, boxstyle="round,pad=0.06",
                 facecolor=DC, edgecolor=DBR, lw=1))
    ax.add_patch(Rectangle((2.6, 2.8), 9, 0.3, facecolor=DBR))
    ax.text(2.8, 2.87, 'All Invoices', fontsize=8, fontweight='bold', color='white')
    p = os.path.join(IMG, "ss_invoices.png")
    plt.savefig(p, dpi=180, bbox_inches='tight', facecolor=DB); plt.close(); return p


def gen_ss_stocks():
    fig, ax = plt.subplots(figsize=(12, 7))
    _sidebar(ax, 5)
    ax.text(2.6, 7.25, 'Stock Management', fontsize=11, fontweight='bold', color='white')
    ax.add_patch(FancyBboxPatch((2.6, 5.8), 9, 1.2, boxstyle="round,pad=0.06",
                 facecolor=DC, edgecolor=DCY, lw=2))
    ax.text(2.8, 6.65, '+ Restock Product', fontsize=8, fontweight='bold', color=DCY)
    kpis = [('50', 'Products', DBL), ('38', 'In Stock', DG),
            ('8', 'Low Stock', DY), ('4', 'Out', DR)]
    for i, (val, label, color) in enumerate(kpis):
        x = 2.6 + i * 2.3
        ax.add_patch(FancyBboxPatch((x, 4.5), 2, 1, boxstyle="round,pad=0.06",
                     facecolor=DC, edgecolor=DBR, lw=1))
        ax.add_patch(Rectangle((x, 5.47), 2, 0.03, facecolor=color))
        ax.text(x + 1, 5.1, val, ha='center', fontsize=16, fontweight='bold', color='white')
        ax.text(x + 1, 4.7, label, ha='center', fontsize=7, color=DM)
    ax.add_patch(FancyBboxPatch((2.6, 0.3), 9, 3.8, boxstyle="round,pad=0.06",
                 facecolor=DC, edgecolor=DBR, lw=1))
    ax.add_patch(Rectangle((2.6, 3.8), 9, 0.3, facecolor=DBR))
    ax.text(2.8, 3.87, 'Stock Overview', fontsize=8, fontweight='bold', color='white')
    p = os.path.join(IMG, "ss_stocks.png")
    plt.savefig(p, dpi=180, bbox_inches='tight', facecolor=DB); plt.close(); return p


def gen_ss_reports():
    fig, ax = plt.subplots(figsize=(12, 7))
    _sidebar(ax, 6)
    ax.text(2.6, 7.25, 'Sales Reports', fontsize=11, fontweight='bold', color='white')
    ax.add_patch(FancyBboxPatch((2.6, 6), 9, 0.8, boxstyle="round,pad=0.06",
                 facecolor=DC, edgecolor=DBL, lw=1))
    for i, f in enumerate(['From', 'To', 'Category', 'Region']):
        ax.add_patch(FancyBboxPatch((2.8 + i * 2.2, 6.2), 2, 0.4,
                     boxstyle="round,pad=0.04", facecolor=DB, edgecolor=DBR, lw=1))
        ax.text(3.8 + i * 2.2, 6.4, f, ha='center', fontsize=6, color=DM)
    kpis = [('$125K', 'Revenue'), ('342', 'Sales'), ('$365', 'Avg')]
    for i, (val, label) in enumerate(kpis):
        x = 2.6 + i * 3.1
        ax.add_patch(FancyBboxPatch((x, 4.7), 2.9, 1, boxstyle="round,pad=0.06",
                     facecolor=DC, edgecolor=DBR, lw=1))
        ax.text(x + 1.45, 5.3, val, ha='center', fontsize=14,
                fontweight='bold', color=DP)
        ax.text(x + 1.45, 4.9, label, ha='center', fontsize=7, color=DM)
    ax.add_patch(FancyBboxPatch((2.6, 0.5), 9, 3.8, boxstyle="round,pad=0.06",
                 facecolor=DC, edgecolor=DBR, lw=1))
    np.random.seed(42)
    x_data = np.linspace(3, 11, 30)
    y_data = 1.5 + 0.5 * np.sin(x_data) + 0.1 * np.random.randn(30)
    ax.plot(x_data, y_data, color=DP, lw=2, marker='o', markersize=2)
    p = os.path.join(IMG, "ss_reports.png")
    plt.savefig(p, dpi=180, bbox_inches='tight', facecolor=DB); plt.close(); return p


# ══════════════════════════════════════════
#  PPTX HELPERS
# ══════════════════════════════════════════

def set_bg(s, c=BG):
    f = s.background.fill; f.solid(); f.fore_color.rgb = c

def add_box(s, l, t, w, h, fc, ec=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fc
    if ec: sh.line.color.rgb = ec; sh.line.width = Pt(1)
    else: sh.line.fill.background()
    sh.adjustments[0] = 0.04; return sh

def add_bar(s, l, t, w, h, c):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    return sh

def add_txt(s, l, t, w, h, text, sz=18, color=BODY_C, bold=False,
            align=PP_ALIGN.LEFT, font=FONT):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
    p.font.color.rgb = color; p.font.bold = bold; p.font.name = font
    p.alignment = align; return tb

def add_img(s, path, l, t, w=None):
    if w: s.shapes.add_picture(path, l, t, width=w)
    else: s.shapes.add_picture(path, l, t)

def sn(s, n):
    add_txt(s, Inches(8.5), Inches(6.85), Inches(1.2), Inches(0.3),
            f"{n}/15", sz=10, color=MUTED, align=PP_ALIGN.RIGHT)

def hdr(s, title, subtitle="", n=1):
    """Slide header — white theme, Times New Roman."""
    add_bar(s, Inches(0), Inches(0), Inches(10), Pt(4), ACCENT)
    add_txt(s, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
            title, sz=36, color=TITLE_C, bold=True)
    if subtitle:
        add_txt(s, Inches(0.6), Inches(1.0), Inches(8.8), Inches(0.4),
                subtitle, sz=16, color=MUTED)
    add_bar(s, Inches(0.6), Inches(1.5), Inches(1.8), Pt(2), ACCENT)
    sn(s, n)


# ══════════════════════════════════════════
#  BUILD 15 SLIDES
# ══════════════════════════════════════════

def build(img):
    prs = Presentation()
    prs.slide_width = Inches(10); prs.slide_height = Inches(7.5)
    bl = prs.slide_layouts[6]

    # ── 1: TITLE PAGE ───────────────────
    s = prs.slides.add_slide(bl); set_bg(s)
    add_bar(s, Inches(0), Inches(0), Inches(10), Pt(5), ACCENT)
    add_bar(s, Inches(0), Inches(7.45), Inches(10), Pt(5), BLUE)
    add_txt(s, Inches(1), Inches(1), Inches(8), Inches(0.8), "SALESDB",
            sz=48, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_txt(s, Inches(1), Inches(1.9), Inches(8), Inches(0.5),
            "Sales Analytics & Invoice Management System",
            sz=22, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_bar(s, Inches(3.5), Inches(2.7), Inches(3), Pt(2), ACCENT)
    add_txt(s, Inches(1), Inches(2.9), Inches(8), Inches(0.4),
            "A Full-Stack Web Application with Real-Time Analytics",
            sz=14, color=MUTED, align=PP_ALIGN.CENTER)
    add_box(s, Inches(2.5), Inches(3.8), Inches(5), Inches(2.8), LIGHT, BORDER)
    info = [
        "Presented by:  [Your Name]",
        "Roll No:  [Your Roll Number]",
        "Department:  [Your Department]",
        "Guide:  Prof. [Guide Name]",
        "Institution:  [Your Institution]",
        "Academic Year:  2024-25",
    ]
    for i, line in enumerate(info):
        add_txt(s, Inches(2.7), Inches(3.95 + i * 0.4), Inches(4.6), Inches(0.35),
                line, sz=13, color=BODY_C, align=PP_ALIGN.CENTER)
    sn(s, 1)

    # ── 2: CONTENT PAGE ────────────────
    s = prs.slides.add_slide(bl); set_bg(s)
    hdr(s, "Content Page", "Table of Contents", 2)
    toc = [
        ("1.", "Introduction", "3"),
        ("2.", "Objective of Work", "4"),
        ("3.", "Technologies Used (Frontend & Backend)", "5"),
        ("4.", "Specification Requirement", "6"),
        ("5.", "Project Flow Diagram — ER Diagram", "7"),
        ("6.", "Project Flow Diagram — DFD", "8"),
        ("7.", "Page Layout (Screenshots — Part 1)", "9"),
        ("8.", "Page Layout (Screenshots — Part 2)", "10"),
        ("9.", "Security Applied", "11"),
        ("10.", "Result Analysis", "12"),
        ("11.", "Conclusion", "13"),
        ("12.", "Future Enhancement", "14"),
        ("13.", "Bibliography / Reference", "15"),
    ]
    for i, (num, title, pg) in enumerate(toc):
        y = Inches(1.8 + i * 0.38)
        clr = ACCENT if i % 2 == 0 else BLUE
        add_txt(s, Inches(0.8), y, Inches(0.6), Inches(0.35),
                num, sz=14, color=clr, bold=True)
        add_txt(s, Inches(1.4), y, Inches(6.5), Inches(0.35),
                title, sz=14, color=BODY_C)
        add_txt(s, Inches(8.5), y, Inches(0.8), Inches(0.35),
                pg, sz=14, color=MUTED, align=PP_ALIGN.RIGHT)
        if i < len(toc) - 1:
            add_bar(s, Inches(1.4), y + Inches(0.32), Inches(7.8), Pt(0.5), BORDER)

    # ── 3: INTRODUCTION ────────────────
    s = prs.slides.add_slide(bl); set_bg(s)
    hdr(s, "Introduction", "About SalesDB", 3)

    intro_points = [
        "SalesDB is a full-stack web application designed for managing sales operations, "
        "customer relationships, product inventory, and invoice generation.",
        "Built using Python Flask as backend, PostgreSQL as database, and a modern "
        "HTML/CSS/JavaScript frontend with Chart.js for interactive analytics.",
        "The system provides a real-time dashboard with KPI cards, revenue charts, "
        "and customer segmentation analysis.",
        "It includes 7 modules: Dashboard, Sales, Customers, Products, "
        "Invoices, Stocks, and Reports.",
        "Admin authentication with session-based login ensures data security "
        "while keeping the dashboard publicly accessible.",
    ]
    for i, point in enumerate(intro_points):
        y = Inches(1.8 + i * 0.95)
        clr = [ACCENT, BLUE, GREEN, PINK, ORANGE][i]
        add_bar(s, Inches(0.6), y + Inches(0.08), Pt(4), Inches(0.3), clr)
        add_txt(s, Inches(0.85), y, Inches(8.8), Inches(0.85),
                point, sz=15, color=BODY_C)

    # ── 4: OBJECTIVE OF WORK ───────────
    s = prs.slides.add_slide(bl); set_bg(s)
    hdr(s, "Objective of Work", "Goals of the Project", 4)

    objs = [
        "Design and develop a full-stack web application for sales data management.",
        "Implement CRUD operations for Customers, Products, and Sales records.",
        "Create an interactive dashboard with KPIs and Chart.js visualizations.",
        "Design a star-schema database with 5 tables (3 dimension + 2 fact).",
        "Develop 20+ RESTful API endpoints using Python Flask.",
        "Implement session-based admin authentication with @login_required decorator.",
        "Build a professional invoice generator with PDF export (ReportLab).",
        "Create real-time stock tracking with low-stock alerts and history logging.",
        "Develop advanced sales reports with date, category, and region filters.",
        "Build a responsive, modern dark-themed UI with glassmorphism effects.",
    ]
    for i, obj in enumerate(objs):
        y = Inches(1.75 + i * 0.52)
        clr = [ACCENT, BLUE, GREEN, PINK, ORANGE, CYAN, ACCENT, BLUE, GREEN, PINK][i]
        add_txt(s, Inches(0.6), y, Inches(0.5), Inches(0.4),
                f"{i + 1:02d}.", sz=13, color=clr, bold=True)
        add_txt(s, Inches(1.2), y, Inches(8.5), Inches(0.45),
                obj, sz=13, color=BODY_C)

    # ── 5: TECHNOLOGIES USED ──────────
    s = prs.slides.add_slide(bl); set_bg(s)
    hdr(s, "Technologies Used", "Frontend & Backend Stack", 5)
    add_img(s, img['tech_stack'], Inches(0.4), Inches(1.8), w=Inches(9.2))

    # Frontend card
    add_box(s, Inches(0.5), Inches(5.2), Inches(4.3), Inches(2), LIGHT, BLUE)
    add_bar(s, Inches(0.5), Inches(5.2), Inches(4.3), Pt(3), BLUE)
    add_txt(s, Inches(0.7), Inches(5.3), Inches(3.9), Inches(0.3),
            "FRONTEND", sz=12, color=BLUE, bold=True)
    fe_items = ["HTML5 — Structure & Semantics", "CSS3 — Dark Theme + Glassmorphism",
                "JavaScript ES6 — DOM + Fetch API", "Chart.js 4.x — Interactive Charts"]
    for i, item in enumerate(fe_items):
        add_txt(s, Inches(0.7), Inches(5.7 + i * 0.32), Inches(3.9), Inches(0.3),
                f"•  {item}", sz=11, color=BODY_C)

    # Backend card
    add_box(s, Inches(5.2), Inches(5.2), Inches(4.3), Inches(2), LIGHT, ACCENT)
    add_bar(s, Inches(5.2), Inches(5.2), Inches(4.3), Pt(3), ACCENT)
    add_txt(s, Inches(5.4), Inches(5.3), Inches(3.9), Inches(0.3),
            "BACKEND", sz=12, color=ACCENT, bold=True)
    be_items = ["Python 3.12 — Core Language", "Flask 3.0 — Web Framework + API",
                "PostgreSQL 16 — Star Schema DB", "ReportLab — PDF Invoice Generation"]
    for i, item in enumerate(be_items):
        add_txt(s, Inches(5.4), Inches(5.7 + i * 0.32), Inches(3.9), Inches(0.3),
                f"•  {item}", sz=11, color=BODY_C)

    # ── 6: SPECIFICATION REQUIREMENT ───
    s = prs.slides.add_slide(bl); set_bg(s)
    hdr(s, "Specification Requirement", "Hardware & Software", 6)

    # Hardware
    add_box(s, Inches(0.5), Inches(1.9), Inches(4.3), Inches(4.8), LIGHT, BLUE)
    add_bar(s, Inches(0.5), Inches(1.9), Inches(4.3), Pt(4), BLUE)
    add_txt(s, Inches(0.7), Inches(2.0), Inches(3.9), Inches(0.35),
            "⚙️  HARDWARE REQUIREMENTS", sz=14, color=BLUE, bold=True)
    hw = [("Processor", "Intel Core i3 or above"),
          ("RAM", "4 GB minimum (8 GB recommended)"),
          ("Storage", "500 MB free disk space"),
          ("Display", "1366 × 768 or higher resolution"),
          ("Network", "Internet connection for fonts/CDN"),
          ("Architecture", "64-bit operating system")]
    for i, (key, val) in enumerate(hw):
        y = Inches(2.6 + i * 0.6)
        add_txt(s, Inches(0.7), y, Inches(1.5), Inches(0.25),
                key, sz=12, color=ACCENT, bold=True)
        add_txt(s, Inches(0.7), y + Inches(0.25), Inches(3.9), Inches(0.3),
                val, sz=11, color=BODY_C)

    # Software
    add_box(s, Inches(5.2), Inches(1.9), Inches(4.3), Inches(4.8), LIGHT, ACCENT)
    add_bar(s, Inches(5.2), Inches(1.9), Inches(4.3), Pt(4), ACCENT)
    add_txt(s, Inches(5.4), Inches(2.0), Inches(3.9), Inches(0.35),
            "💻  SOFTWARE REQUIREMENTS", sz=14, color=ACCENT, bold=True)
    sw = [("OS", "Windows 10/11, Ubuntu 20+, macOS"),
          ("Python", "Version 3.8 or above"),
          ("PostgreSQL", "Version 14 or above"),
          ("Browser", "Chrome, Firefox, Edge (latest)"),
          ("pip Packages", "flask, psycopg2, reportlab, flask-cors"),
          ("IDE", "VS Code, PyCharm (optional)")]
    for i, (key, val) in enumerate(sw):
        y = Inches(2.6 + i * 0.6)
        add_txt(s, Inches(5.4), y, Inches(1.5), Inches(0.25),
                key, sz=12, color=ACCENT, bold=True)
        add_txt(s, Inches(5.4), y + Inches(0.25), Inches(3.9), Inches(0.3),
                val, sz=11, color=BODY_C)

    # ── 7: ER DIAGRAM ─────────────────
    s = prs.slides.add_slide(bl); set_bg(s)
    hdr(s, "Project Flow — ER Diagram", "5-Table Star Schema Design", 7)
    add_img(s, img['er_diagram'], Inches(0.1), Inches(1.65), w=Inches(9.8))

    # ── 8: DFD ─────────────────────────
    s = prs.slides.add_slide(bl); set_bg(s)
    hdr(s, "Project Flow — DFD", "Data Flow Diagrams (Level 0 & Level 1)", 8)
    add_txt(s, Inches(0.4), Inches(1.7), Inches(4.8), Inches(0.3),
            "DFD Level 0 (Context Diagram)", sz=14, color=ACCENT, bold=True)
    add_img(s, img['dfd_level0'], Inches(0.1), Inches(2.05), w=Inches(4.9))
    add_txt(s, Inches(5.2), Inches(1.7), Inches(4.8), Inches(0.3),
            "DFD Level 1 (Detailed Processes)", sz=14, color=BLUE, bold=True)
    add_img(s, img['dfd_level1'], Inches(5), Inches(2.05), w=Inches(4.9))

    # ── 9: PAGE LAYOUT — SCREENSHOTS 1 ──
    s = prs.slides.add_slide(bl); set_bg(s)
    hdr(s, "Page Layout — Screenshots", "Dashboard, Login & Sales Pages", 9)

    add_txt(s, Inches(0.4), Inches(1.7), Inches(3), Inches(0.3),
            "Dashboard (Public)", sz=13, color=ACCENT, bold=True)
    add_img(s, img['ss_dashboard'], Inches(0.3), Inches(2.0), w=Inches(3.1))

    add_txt(s, Inches(3.6), Inches(1.7), Inches(3), Inches(0.3),
            "Admin Login", sz=13, color=ORANGE, bold=True)
    add_img(s, img['ss_login'], Inches(3.5), Inches(2.0), w=Inches(3.1))

    add_txt(s, Inches(6.8), Inches(1.7), Inches(3), Inches(0.3),
            "Sales Management", sz=13, color=BLUE, bold=True)
    add_img(s, img['ss_sales'], Inches(6.7), Inches(2.0), w=Inches(3.1))

    # Bottom row
    add_txt(s, Inches(0.4), Inches(5.0), Inches(3), Inches(0.3),
            "Products", sz=13, color=GREEN, bold=True)
    add_img(s, img['ss_products'], Inches(0.3), Inches(5.3), w=Inches(3.1))

    # ── 10: PAGE LAYOUT — SCREENSHOTS 2 ──
    s = prs.slides.add_slide(bl); set_bg(s)
    hdr(s, "Page Layout — Screenshots", "Invoices, Stocks & Reports", 10)

    add_txt(s, Inches(0.4), Inches(1.7), Inches(3), Inches(0.3),
            "Invoice Builder", sz=13, color=ACCENT, bold=True)
    add_img(s, img['ss_invoices'], Inches(0.3), Inches(2.0), w=Inches(3.1))

    add_txt(s, Inches(3.6), Inches(1.7), Inches(3), Inches(0.3),
            "Stock Management", sz=13, color=CYAN, bold=True)
    add_img(s, img['ss_stocks'], Inches(3.5), Inches(2.0), w=Inches(3.1))

    add_txt(s, Inches(6.8), Inches(1.7), Inches(3), Inches(0.3),
            "Reports & Analytics", sz=13, color=PINK, bold=True)
    add_img(s, img['ss_reports'], Inches(6.7), Inches(2.0), w=Inches(3.1))

    # ── 11: SECURITY APPLIED ──────────
    s = prs.slides.add_slide(bl); set_bg(s)
    hdr(s, "Security Applied", "Authentication & Access Control", 11)

    add_img(s, img['auth_flow'], Inches(0.2), Inches(1.7), w=Inches(5.5))

    add_box(s, Inches(5.9), Inches(1.7), Inches(3.8), Inches(5.5), LIGHT, ACCENT)
    add_bar(s, Inches(5.9), Inches(1.7), Inches(3.8), Pt(4), ACCENT)
    add_txt(s, Inches(6.1), Inches(1.8), Inches(3.4), Inches(0.3),
            "🔐  SECURITY FEATURES", sz=13, color=ACCENT, bold=True)

    sec_features = [
        ("Session-Based Auth", "Flask sessions with secure cookies"),
        ("30-Min Expiry", "Auto-logout after inactivity"),
        ("@login_required", "Python decorator protects routes"),
        ("Public Dashboard", "GET endpoints remain accessible"),
        ("Protected CRUD", "POST/PUT/DELETE require login"),
        ("401 Auto-Redirect", "Expired session → login overlay"),
        ("CORS Protection", "flask-cors for cross-origin security"),
        ("Password Validation", "Server-side credential checking"),
    ]
    for i, (title, desc) in enumerate(sec_features):
        y = Inches(2.3 + i * 0.6)
        add_txt(s, Inches(6.1), y, Inches(3.4), Inches(0.25),
                f"▸  {title}", sz=11, color=TITLE_C, bold=True)
        add_txt(s, Inches(6.3), y + Inches(0.25), Inches(3.2), Inches(0.25),
                desc, sz=10, color=MUTED)

    # ── 12: RESULT ANALYSIS ───────────
    s = prs.slides.add_slide(bl); set_bg(s)
    hdr(s, "Result Analysis", "Dashboard Analytics & KPIs", 12)

    add_img(s, img['result_charts'], Inches(0.3), Inches(1.7), w=Inches(9.4))

    # KPI summary cards
    kpis = [
        ("$100,000", "Total Revenue", ACCENT),
        ("245", "Total Sales", BLUE),
        ("$408.16", "Avg Order Value", GREEN),
        ("48", "Total Customers", PINK),
        ("12", "Total Products", ORANGE),
    ]
    for i, (val, label, clr) in enumerate(kpis):
        x = Inches(0.3 + i * 1.9)
        add_box(s, x, Inches(5.5), Inches(1.75), Inches(1.4), LIGHT, clr)
        add_bar(s, x, Inches(5.5), Inches(1.75), Pt(3), clr)
        add_txt(s, x + Inches(0.1), Inches(5.65), Inches(1.55), Inches(0.5),
                val, sz=18, color=clr, bold=True, align=PP_ALIGN.CENTER)
        add_txt(s, x + Inches(0.1), Inches(6.2), Inches(1.55), Inches(0.5),
                label, sz=10, color=MUTED, align=PP_ALIGN.CENTER)

    # ── 13: CONCLUSION ─────────────────
    s = prs.slides.add_slide(bl); set_bg(s)
    hdr(s, "Conclusion", "Project Summary", 13)

    add_box(s, Inches(0.5), Inches(1.8), Inches(9), Inches(1), LIGHT, GREEN)
    add_bar(s, Inches(0.5), Inches(1.8), Inches(9), Pt(4), GREEN)
    add_txt(s, Inches(0.7), Inches(1.9), Inches(8.6), Inches(0.3),
            "✅  ALL PROJECT OBJECTIVES SUCCESSFULLY ACHIEVED", sz=15, color=GREEN, bold=True)
    add_txt(s, Inches(0.7), Inches(2.3), Inches(8.6), Inches(0.4),
            "SalesDB delivers a complete, functional sales management system with modern UI.",
            sz=13, color=BODY_C)

    conclusions = [
        "Developed a full-stack web application using Flask, PostgreSQL, and modern JavaScript.",
        "Implemented comprehensive CRUD operations across 7 application modules.",
        "Created 5 relational database tables following star schema design.",
        "Built 20+ RESTful API endpoints with session-based authentication.",
        "Delivered professional PDF invoice generation using ReportLab.",
        "Integrated real-time stock tracking with automated alerts.",
        "Provided advanced reporting with multi-filter support and CSV export.",
    ]
    for i, c in enumerate(conclusions):
        y = Inches(3.1 + i * 0.55)
        clr = [ACCENT, BLUE, GREEN, PINK, CYAN, ORANGE, ACCENT][i]
        add_bar(s, Inches(0.6), y + Inches(0.06), Pt(4), Inches(0.25), clr)
        add_txt(s, Inches(0.85), y, Inches(8.8), Inches(0.5), c, sz=13, color=BODY_C)

    # Stats
    stats = [("7", "Modules", ACCENT), ("20+", "API Endpoints", BLUE),
             ("5", "DB Tables", GREEN), ("15+", "Features", PINK)]
    for i, (val, label, clr) in enumerate(stats):
        x = Inches(0.5 + i * 2.35)
        add_box(s, x, Inches(6.6), Inches(2.15), Inches(0.7), LIGHT, clr)
        add_txt(s, x + Inches(0.1), Inches(6.62), Inches(0.7), Inches(0.6),
                val, sz=18, color=clr, bold=True, align=PP_ALIGN.CENTER)
        add_txt(s, x + Inches(0.8), Inches(6.7), Inches(1.25), Inches(0.5),
                label, sz=11, color=BODY_C)

    # ── 14: FUTURE ENHANCEMENT ─────────
    s = prs.slides.add_slide(bl); set_bg(s)
    hdr(s, "Future Enhancement", "Planned Improvements", 14)

    enhancements = [
        ("💳", "Payment Gateway Integration", "Stripe/Razorpay for online payments", ACCENT),
        ("📧", "Email Invoice Delivery", "Auto-send PDF invoices via SendGrid/SMTP", BLUE),
        ("👥", "Multi-User Role System", "Admin, Manager, Employee with different permissions", GREEN),
        ("📱", "Mobile Application", "React Native / Flutter companion app", PINK),
        ("🔍", "Barcode / QR Scanner", "Quick product lookup and inventory scanning", ORANGE),
        ("☁️", "Cloud Deployment", "Deploy on AWS / Heroku / Railway with CI/CD", CYAN),
        ("📊", "Advanced Analytics", "Predictive sales forecasting with ML models", ACCENT),
        ("🔔", "Real-time Notifications", "WebSocket-based live updates for stock alerts", BLUE),
    ]
    for i, (icon, title, desc, clr) in enumerate(enhancements):
        r = i // 2; c = i % 2
        x = Inches(0.4 + c * 4.8)
        y = Inches(1.8 + r * 1.25)
        add_box(s, x, y, Inches(4.6), Inches(1.05), LIGHT, clr)
        add_bar(s, x, y, Pt(4), Inches(1.05), clr)
        add_txt(s, x + Inches(0.15), y + Inches(0.08), Inches(4.3), Inches(0.35),
                f"{icon}  {title}", sz=13, color=clr, bold=True)
        add_txt(s, x + Inches(0.15), y + Inches(0.5), Inches(4.3), Inches(0.5),
                desc, sz=11, color=MUTED)

    # ── 15: BIBLIOGRAPHY ───────────────
    s = prs.slides.add_slide(bl); set_bg(s)
    hdr(s, "Bibliography / Reference", "Resources & Documentation", 15)

    refs = [
        ("Flask Documentation", "https://flask.palletsprojects.com/", ACCENT),
        ("PostgreSQL Documentation", "https://www.postgresql.org/docs/", BLUE),
        ("Chart.js Documentation", "https://www.chartjs.org/docs/", GREEN),
        ("ReportLab User Guide", "https://www.reportlab.com/docs/", PINK),
        ("MDN Web Docs (HTML/CSS/JS)", "https://developer.mozilla.org/", ORANGE),
        ("psycopg2 Documentation", "https://www.psycopg.org/docs/", CYAN),
        ("python-pptx Documentation", "https://python-pptx.readthedocs.io/", ACCENT),
        ("W3Schools Web Tutorials", "https://www.w3schools.com/", BLUE),
        ("Stack Overflow Community", "https://stackoverflow.com/", GREEN),
        ("GitHub — Version Control", "https://github.com/", PINK),
    ]
    for i, (title, url, clr) in enumerate(refs):
        y = Inches(1.8 + i * 0.5)
        add_txt(s, Inches(0.6), y, Inches(0.5), Inches(0.4),
                f"[{i + 1}]", sz=12, color=clr, bold=True)
        add_txt(s, Inches(1.2), y, Inches(4.5), Inches(0.4),
                title, sz=13, color=TITLE_C, bold=True)
        add_txt(s, Inches(5.8), y, Inches(3.8), Inches(0.4),
                url, sz=11, color=MUTED)
        if i < len(refs) - 1:
            add_bar(s, Inches(1.2), y + Inches(0.4), Inches(8.2), Pt(0.5), BORDER)

    fname = 'SalesDB_Presentation.pptx'
    prs.save(fname); return fname


# ══════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════

if __name__ == '__main__':
    print('\n  ╔══════════════════════════════════════════╗')
    print('  ║  SalesDB — 15-Slide PPT (White Theme)    ║')
    print('  ║  Font: Times New Roman                    ║')
    print('  ║  Title < 40pt | Body < 26pt               ║')
    print('  ╚══════════════════════════════════════════╝\n')
    print('  → Generating diagrams & screenshots...')

    img = {}

    # White-background diagrams
    img['tech_stack'] = gen_tech_stack();       print('    ✓ Tech Stack')
    img['er_diagram'] = gen_er_diagram();       print('    ✓ ER Diagram (5 tables)')
    img['dfd_level0'] = gen_dfd_level0();       print('    ✓ DFD Level 0')
    img['dfd_level1'] = gen_dfd_level1();       print('    ✓ DFD Level 1 (7 processes)')
    img['auth_flow']  = gen_auth_flow();        print('    ✓ Auth Flow')
    img['result_charts'] = gen_result_charts(); print('    ✓ Result Charts')

    # Dark-background screenshots (actual app look)
    img['ss_dashboard'] = gen_ss_dashboard();   print('    ✓ Screenshot: Dashboard')
    img['ss_login']     = gen_ss_login();       print('    ✓ Screenshot: Login')
    img['ss_sales']     = gen_ss_sales();       print('    ✓ Screenshot: Sales')
    img['ss_products']  = gen_ss_products();    print('    ✓ Screenshot: Products')
    img['ss_invoices']  = gen_ss_invoices();    print('    ✓ Screenshot: Invoices')
    img['ss_stocks']    = gen_ss_stocks();      print('    ✓ Screenshot: Stocks')
    img['ss_reports']   = gen_ss_reports();     print('    ✓ Screenshot: Reports')

    print(f'\n  → 13 images saved to {IMG}/')
    print('  → Building 15 slides...')

    fname = build(img)

    print(f'\n  ✅ Saved: {fname}')
    print('  ✅ Total: 15 slides')
    print('  ✅ Theme: White / Clean')
    print('  ✅ Font:  Times New Roman')
    print('  ✅ Title: 36pt | Body: 13-18pt\n')

    slides = [
        '01. Title / Cover Page',
        '02. Content Page (Table of Contents)',
        '03. Introduction',
        '04. Objective of Work (10 goals)',
        '05. Technologies Used (Frontend & Backend)',
        '06. Specification Requirement (HW & SW)',
        '07. Project Flow — ER Diagram [5 TABLES]',
        '08. Project Flow — DFD [Level 0 + Level 1]',
        '09. Page Layout — Screenshots Part 1 (Dashboard, Login, Sales, Products)',
        '10. Page Layout — Screenshots Part 2 (Invoices, Stocks, Reports)',
        '11. Security Applied (Auth Flow + Features)',
        '12. Result Analysis (Charts + KPIs)',
        '13. Conclusion',
        '14. Future Enhancement (8 items)',
        '15. Bibliography / Reference (10 sources)',
    ]
    print('  Slides:')
    for sl in slides:
        print(f'    {sl}')
    print()